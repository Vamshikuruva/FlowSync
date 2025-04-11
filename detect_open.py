import psutil
import win32gui
import win32process
import os
import time
import shutil
import tempfile
import pygetwindow as gw
import pyperclip
import keyboard
import subprocess
import docx
import PyPDF2
import pandas as pd
import pptx
import json
import pyautogui
import pickle
from dotenv import load_dotenv
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain_openai import OpenAIEmbeddings
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain.schema.output_parser import StrOutputParser
from langchain.schema.runnable import RunnablePassthrough
from langchain_community.vectorstores import FAISS
from urllib.parse import unquote

TEMP_DIR = tempfile.gettempdir()  

FILE_TYPES = {
    '.docx': "Word Document",
    '.pdf': "PDF Document",
    '.txt': "Text File",
    '.xlsx': "Excel Spreadsheet",
    '.xls': "Excel Spreadsheet",
    '.pptx': "PowerPoint Presentation",
    '.ppt': "PowerPoint Presentation",
    '.csv': "CSV File",
    '.json': "JSON File"
}

BROWSER_PROCESSES = ["chrome.exe", "msedge.exe", "firefox.exe"]

def get_active_window():
    return win32gui.GetForegroundWindow()

def get_process_from_window(hwnd):
    _, pid = win32process.GetWindowThreadProcessId(hwnd)
    return psutil.Process(pid)

def get_open_file_path(process):
    try:
        for file in process.open_files():
            if any(file.path.endswith(ext) for ext in FILE_TYPES.keys()):
                return file.path
    except (psutil.AccessDenied, Exception) as e:
        print(f"[ERROR] Could not retrieve file path: {e}")
    return None

def close_application_by_pid(pid):
    """Force close an application using its process ID."""
    if pid:
        try:
            print(f"🔴 Closing process (PID: {pid})")
            psutil.Process(pid).terminate()
            time.sleep(2)  
        except psutil.NoSuchProcess:
            print("⚠️ Process already closed.")

def reopen_file(file_path):
    """Reopens the file after saving changes."""
    print(f"🔄 Reopening {file_path}")
    subprocess.Popen(['start', '', file_path], shell=True)

def copy_to_temp(file_path):
    """Copy the file to a temporary folder for editing."""
    temp_path = os.path.join(TEMP_DIR, os.path.basename(file_path))
    shutil.copy2(file_path, temp_path)
    return temp_path

def capture_screenshot():
    screenshot_path = os.path.join(tempfile.gettempdir(), "screenshot.png")
    screenshot = pyautogui.screenshot()
    screenshot.save(screenshot_path)
    print(f"📸 Screenshot saved: {screenshot_path}")
    return screenshot_path

def get_browser_pdf_url():
    active_window = gw.getActiveWindow()
    if not active_window:
        return None
    keyboard.press_and_release("ctrl+l")
    time.sleep(0.5)
    keyboard.press_and_release("ctrl+c")
    time.sleep(0.5)
    url = pyperclip.paste()
    if url.startswith("file:///"):
        return unquote(url[8:]).replace("/", "\\")
    return None

def detect_document_path():
    time.sleep(3)  
    hwnd = get_active_window()
    process = get_process_from_window(hwnd)
    file_path = get_open_file_path(process)
    if not file_path:
        if process.name().lower() in BROWSER_PROCESSES:
            file_path = get_browser_pdf_url()
    return file_path, process

def extract_text(file_path):
    """Extracts text from various file formats."""
    if not os.path.exists(file_path):
        return "File not found."
    file_extension = file_path.lower().split('.')[-1]
    try:
        if file_extension == 'txt':
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        elif file_extension == 'docx':
            doc = docx.Document(file_path)
            return '\n'.join([para.text for para in doc.paragraphs])
        elif file_extension == 'pdf':
            text = ""
            with open(file_path, "rb") as f:
                reader = PyPDF2.PdfReader(f)
                for page in reader.pages:
                    text += page.extract_text() + "\n"
            return text if text else "No text found in PDF."
        elif file_extension in ['xlsx', 'xls']:
            df = pd.read_excel(file_path, sheet_name=None)
            text = [f"--- Sheet: {sheet} ---\n{data.to_string()}" for sheet, data in df.items()]
            return "\n".join(text)
        elif file_extension in ['pptx', 'ppt']:
            ppt = pptx.Presentation(file_path)
            return "\n".join([shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")])
        elif file_extension == 'csv':
            df = pd.read_csv(file_path)
            return df.to_string(index=False)
        elif file_extension == 'json':
            with open(file_path, 'r', encoding='utf-8') as f:
                return json.dumps(json.load(f), indent=4)
        else:
            return "Unsupported file type."
    except Exception as e:
        return f"Error extracting text: {e}"

def data_chunks(data):
    """Splits data into chunks for embedding."""
    return RecursiveCharacterTextSplitter(chunk_size=1000, chunk_overlap=200, length_function=len).split_text(data)

def chunk_embedding(chunks, file_name, openai_api_key):
    """Embeds the chunks using FAISS and Google embeddings."""
    embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=openai_api_key)
    index_path = f"{file_name}_index"
    if os.path.exists(index_path):
        VectorStore = FAISS.load_local(index_path, embeddings, allow_dangerous_deserialization=True)
    else:
        VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
        VectorStore.save_local(index_path)
    return VectorStore

def load_permanent_index(index_path, openai_api_key):
    embeddings = OpenAIEmbeddings(model="text-embedding-3-large", openai_api_key=openai_api_key)
    if os.path.exists(index_path):
        return FAISS.load_local(index_path, embeddings, allow_dangerous_deserialization=True)
    else:
        print("⚠️ Permanent index not found.")
        return None

def build_temp_index_from_file(file_path, openai_api_key):
    embeddings = OpenAIEmbeddings(
        model="text-embedding-3-large",
        openai_api_key=openai_api_key
    )

    base_name = os.path.basename(file_path)[:-4]
    temp_index_path = os.path.join(TEMP_DIR, f"{base_name}_temp_index")
    index_path = f"{base_name}_index"  # Assuming saved in current working directory
    # Step 1: Check TEMP_DIR for temp index
    if os.path.exists(temp_index_path):
        print(f"📦 Found existing TEMP index: {temp_index_path}")
        return FAISS.load_local(temp_index_path, embeddings, allow_dangerous_deserialization=True)
    # Step 2: Check for permanent index in current dir
    if os.path.exists(index_path):
        print(f"📦 Found existing index: {index_path}")
        return FAISS.load_local(index_path, embeddings, allow_dangerous_deserialization=True)
    # Step 3: Build a new temp index if none found
    print("🧠 No existing index found. Creating new TEMP index...")
    data = extract_text(file_path)
    if data:
        chunks = data_chunks(data)
        VectorStore = FAISS.from_texts(chunks, embedding=embeddings)
        VectorStore.save_local(temp_index_path)
        print(f"💾 Index saved at: {temp_index_path}")
        return VectorStore
    print("⚠️ Failed to extract or embed content from the file.")
    return None


if __name__ == "__main__":
    file_path, process = detect_document_path()
    print("file_path:", file_path)
    print("process:", process.name())
    
    if file_path:
        print("📄 Opened file:", file_path)
        load_dotenv()
        openai_api_key = os.getenv("OPENAI_API_KEY")

        print("🔴 Closing application...")
        close_application_by_pid(process.pid)
        print("✅ Application closed.")

        try:
            temp_path = copy_to_temp(file_path)
            print(f"✅ File copied to temp: {temp_path}")
        except PermissionError as e:
            print(f"❌ Permission denied: {e}")

        print("🔄 Reopening file...")
        reopen_file(file_path)

        # Load permanent index
        permanent_index = load_permanent_index("permanent_index", openai_api_key)
        
        # Build temporary index from opened file
        temp_index = build_temp_index_from_file(file_path, openai_api_key)

        if not temp_index and not permanent_index:
            print("❌ No documents available to answer from.")
            exit()

        # Combine documents for similarity search
        combined_docs = []
        while True:
            query = input("Enter a query or type EXIT: ")
            if query.lower() == "exit":
                break

            if temp_index:
                combined_docs += temp_index.similarity_search(query, k=3)
            if permanent_index:
                combined_docs += permanent_index.similarity_search(query, k=3)

            if not combined_docs:
                print("🤖 No relevant answers found.")
                continue

            context = "\n".join([doc.page_content for doc in combined_docs])
            llm = ChatOpenAI(model='gpt-3.5-turbo', temperature=0.2, api_key=openai_api_key)
            prompt = ChatPromptTemplate.from_template("Answer the question based on: {context} Question: {question}")
            chain = {"context": lambda x: context, "question": RunnablePassthrough()} | prompt | llm | StrOutputParser()
            print(chain.invoke(query))
    else:
        print("❌ No document detected. Taking a screenshot...")
        capture_screenshot()
